# -*- coding: utf-8 -*-
"""Build the Project A presentation deck (.pptx) with python-pptx.
Matches the product report + Friday demo script. 16:9 widescreen."""
import os, csv
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "07_paper_draft", "figures")
QUAL = os.path.join(ROOT, "06_method_results", "Qualitative")
DASH = os.path.join(ROOT, "09_dashboard")
OUT = os.path.join(ROOT, "07_paper_draft", "Project_A_LOCO_Presentation.pptx")

# palette
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
DEEP   = RGBColor(0x12, 0x24, 0x3D)
TEAL   = RGBColor(0x1A, 0xB3, 0xA3)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)
PANEL  = RGBColor(0xF1, 0xF5, 0xFA)
CARD   = RGBColor(0xE7, 0xEF, 0xF8)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x2B, 0x35)
GREY   = RGBColor(0x5A, 0x67, 0x74)
HF = "Trebuchet MS"   # header font
BF = "Calibri"        # body font

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    return tb, tf

def para(tf, text, size, color, bold=False, font=BF, align=PP_ALIGN.LEFT,
         first=False, space_after=6, italic=False, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.name = font
    f.color.rgb = color
    if bullet:
        _set_bullet(p)
    return p

def _set_bullet(p):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    bu = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
    pPr.append(bu)

def rect(s, l, t, w, h, fill, line=None, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp

def title(s, text, color=NAVY):
    _, tf = box(s, 0.6, 0.42, 12.1, 1.0)
    para(tf, text, 30, color, bold=True, font=HF, first=True, space_after=0)

def img_fit(s, path, l, t, max_w, max_h, center=True):
    if not os.path.exists(path):
        return None
    iw, ih = Image.open(path).size
    ar = iw / ih
    w, h = max_w, max_w / ar
    if h > max_h:
        h, w = max_h, max_h * ar
    if center:
        l = l + (max_w - w) / 2
        t = t + (max_h - h) / 2
    return s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def first_existing(*paths):
    for p in paths:
        if os.path.exists(p):
            return p
    return paths[-1]

# results
mby = {}
with open(os.path.join(DASH, "corrected_main_results.csv"), newline="") as f:
    for r in csv.DictReader(f):
        mby[r["method"]] = r

# ---------------- S1 Title ----------------
s = slide(DEEP)
rect(s, 0.0, 0.0, 0.28, SH, TEAL, rounded=False)   # thin side motif (not under title)
_, tf = box(s, 1.0, 2.35, 11.3, 2.2)
para(tf, "LOCO Anomaly Inspector", 46, WHITE, bold=True, font=HF, first=True, space_after=8)
para(tf, "An interactive, reproducible system for logical & structural anomaly detection",
     21, ICE, font=HF, space_after=2)
para(tf, "on the MVTec LOCO AD industrial benchmark", 21, ICE, font=HF, space_after=20)
para(tf, "Data Analytics Lab — Project A   •   Final Project Presentation", 14, TEAL, bold=True)

# ---------------- S2 Problem ----------------
s = slide(WHITE)
title(s, "Two kinds of defect — one is much harder")
_, tf = box(s, 0.6, 1.5, 6.0, 5.3)
para(tf, "Structural anomalies", 19, NAVY, bold=True, font=HF, first=True, space_after=2)
para(tf, "Local appearance defects — scratches, dents, cracks. Visible in a small region.",
     15, DARK, space_after=14)
para(tf, "Logical anomalies", 19, NAVY, bold=True, font=HF, space_after=2)
para(tf, "Global rule violations — a missing fruit, an extra pushpin, a swapped part. "
        "Every region looks normal; only the overall composition is wrong.", 15, DARK,
        space_after=14)
para(tf, "Trained only on defect-free images — no labelled defects.", 15, TEAL, bold=True,
        italic=True)
img_fit(s, os.path.join(FIG, "figure_2_dataset_examples.png"), 6.9, 1.55, 5.9, 5.1)

# ---------------- S3 Objective ----------------
s = slide(PANEL)
title(s, "What we built — and why")
_, tf = box(s, 0.6, 1.45, 12.1, 0.7)
para(tf, "A self-contained dashboard + reproducible ML pipeline that detects and explains "
        "anomalies in industrial product images.", 16, GREY, italic=True, first=True)
goals = [
    ("Detect both types", "Flag logical and structural defects from good images only."),
    ("Explain every decision", "Show the anomaly score, the threshold, and where the defect is."),
    ("Relevant analysis only", "Surface the EDA that informs the model — not decorative charts."),
    ("Fully reproducible", "Fixed seeds, leakage-checked data, saved configs — every result is auditable."),
]
gx = [0.6, 6.86]; gy = [2.35, 4.65]; cw, ch = 5.87, 2.05
for i, (h, d) in enumerate(goals):
    l = gx[i % 2]; t = gy[i // 2]
    rect(s, l, t, cw, ch, WHITE)
    _, tf = box(s, l + 0.3, t + 0.25, cw - 0.6, ch - 0.5)
    para(tf, h, 18, NAVY, bold=True, font=HF, first=True, space_after=6)
    para(tf, d, 14.5, DARK)

# ---------------- S4 System overview ----------------
s = slide(WHITE)
title(s, "How it works — end to end")
img_fit(s, os.path.join(FIG, "figure_1_pipeline_overview.png"), 0.7, 1.5, 12.0, 4.6)
_, tf = box(s, 0.6, 6.35, 12.1, 0.8)
para(tf, "Data audit & letterbox preprocessing → EDA → four detectors → validation-only "
        "fusion → evaluation → dashboard & report. All logic in one shared module.",
        14, GREY, italic=True, first=True, align=PP_ALIGN.CENTER)

# ---------------- S5 Data & rigor ----------------
s = slide(DEEP)
title(s, "Data & preprocessing rigor", color=WHITE)
stats = [("3,651", "product images"), ("5", "categories"),
         ("0", "cross-split leaks"), ("384²", "letterbox resize")]
sx = 0.7; cw = 2.92; gap = 0.18
for i, (num, lab) in enumerate(stats):
    l = sx + i * (cw + gap)
    rect(s, l, 1.7, cw, 1.9, NAVY)
    _, tf = box(s, l + 0.1, 1.82, cw - 0.2, 1.7)
    para(tf, num, 40, TEAL, bold=True, font=HF, first=True, align=PP_ALIGN.CENTER, space_after=2)
    para(tf, lab, 14, ICE, align=PP_ALIGN.CENTER)
_, tf = box(s, 0.7, 4.1, 11.9, 2.6)
para(tf, "Raw dataset treated as read-only.", 16, WHITE, bold=True, first=True, space_after=8, bullet=True)
para(tf, "Aspect-preserving letterbox to 384×384 (bilinear images, nearest-neighbour masks).",
     16, WHITE, space_after=8, bullet=True)
para(tf, "MD5 + perceptual hashing across splits to detect duplicate leakage — none found.",
     16, WHITE, space_after=8, bullet=True)
para(tf, "Fit on Train(good); threshold on Validation(good); Test held out. No test labels in tuning.",
     16, WHITE, bullet=True)

# ---------------- S6 Methods ----------------
s = slide(PANEL)
title(s, "Four DINOv2 detectors + fusion")
methods = [
    ("DINOv2 Patch Memory", "Memory bank of normal DINOv2 patch tokens; cosine nearest-neighbour distance. Catches structural defects."),
    ("DINOv2 Region-Aware Memory", "Memory conditioned on image region — adds positional awareness. Strongest single detector."),
    ("DINOv2 Composition Histogram", "DINOv2 visual-word histogram scored by Mahalanobis distance. Targets logical anomalies."),
    ("Global Image-Stat (proxy)", "Single global descriptor vs. the normal mean. Fast image-level baseline (non-DINOv2)."),
]
gx = [0.6, 6.86]; gy = [1.5, 3.55]; cw, ch = 5.87, 1.82
for i, (h, d) in enumerate(methods):
    l = gx[i % 2]; t = gy[i // 2]
    rect(s, l, t, cw, ch, WHITE)
    _, tf = box(s, l + 0.28, t + 0.2, cw - 0.56, ch - 0.4)
    para(tf, h, 17, NAVY, bold=True, font=HF, first=True, space_after=5)
    para(tf, d, 14, DARK)
rect(s, 0.6, 5.65, 12.13, 1.2, NAVY)
_, tf = box(s, 0.9, 5.82, 11.5, 0.9)
para(tf, "Fusion", 17, TEAL, bold=True, font=HF, first=True, space_after=3)
para(tf, "Each detector is z-normalised on Validation(good) only, then combined (mean / max / "
        "rank-average). Complementary detectors fused beat any single one.", 14.5, WHITE)

# ---------------- S7 Results ----------------
s = slide(WHITE)
title(s, "Results — fusion wins")
order = ["Fusion (mean)", "DINOv2 Region-Aware Memory", "DINOv2 Patch Memory (NN)",
         "DINOv2 Composition Histogram (BoVW)", "Global Image-Stat Detector (proxy)"]
short = {"Fusion (mean)": "Fusion (mean)", "DINOv2 Region-Aware Memory": "DINOv2 Region-Aware",
         "DINOv2 Patch Memory (NN)": "DINOv2 Patch Memory",
         "DINOv2 Composition Histogram (BoVW)": "Composition Histogram",
         "Global Image-Stat Detector (proxy)": "Global Image-Stat (proxy)"}
rows = len(order) + 1
tbl_shape = s.shapes.add_table(rows, 4, Inches(0.6), Inches(1.6), Inches(7.6), Inches(3.6))
table = tbl_shape.table
table.columns[0].width = Inches(3.4)
for c in range(1, 4):
    table.columns[c].width = Inches(1.4)
hdrs = ["Method", "Logical", "Structural", "Overall"]
for c, htxt in enumerate(hdrs):
    cell = table.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
    r = p.add_run(); r.text = htxt; r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = HF
for ri, m in enumerate(order, start=1):
    r = mby[m]
    best = (m == "Fusion (mean)")
    vals = [short[m], f"{float(r['logical']):.3f}", f"{float(r['structural']):.3f}", f"{float(r['overall']):.3f}"]
    for c, v in enumerate(vals):
        cell = table.cell(ri, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = CARD if best else WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
        rr = p.add_run(); rr.text = v; rr.font.size = Pt(13); rr.font.bold = best
        rr.font.color.rgb = NAVY if best else DARK; rr.font.name = BF
_, tf = box(s, 8.45, 1.7, 4.3, 4.0)
para(tf, "Image-level AUROC, averaged over 5 categories.", 13, GREY, italic=True, first=True, space_after=14)
para(tf, "Fusing a local detector with a global one beats either alone.", 15.5, DARK, bullet=True, space_after=10)
para(tf, "Region-aware memory is the strongest single detector and leads on structural.", 15.5, DARK, bullet=True, space_after=10)
para(tf, "Real frozen DINOv2-small features — verified end-to-end and fully reproducible.", 15.5, DARK, bullet=True)

# ---------------- S7b Feature selection ----------------
s = slide(PANEL)
title(s, "Feature selection — fewer, better representations")
_, tf = box(s, 0.6, 1.5, 6.0, 5.3)
para(tf, "Each detector is a feature representation of the image. We rank them by single-branch AUROC "
        "and by leave-one-out contribution to the fusion.", 15.5, DARK, first=True, space_after=12)
para(tf, "PatchCore is byte-identical to DINOv2 PatchMemory — a redundant duplicate.", 15.5, DARK, bullet=True, space_after=8)
para(tf, "The Composition-Histogram branch is the only one that hurts the fusion (Δ < 0).", 15.5, DARK, bullet=True, space_after=8)
para(tf, "Dropping both → a leaner 3-representation model: 0.764 overall, 0.826 structural — "
        "edging the full fusion with fewer features.", 16, NAVY, bold=True, bullet=True, space_after=8)
para(tf, "No anomalous validation images exist, so importance is test-assessed and reported as analysis; "
        "the headline 0.76 fusion stays pre-specified and untuned.", 13, GREY, italic=True)
img_fit(s, os.path.join(FIG, "figure_feature_selection.png"), 6.85, 1.7, 6.0, 4.8)

# ---------------- S8 Dashboard / explainability ----------------
s = slide(PANEL)
title(s, "An explainable, offline dashboard")
_, tf = box(s, 0.6, 1.5, 5.5, 5.2)
for h in ["Overview", "EDA (relevant only)", "Leaderboard", "Explainability", "Per-category drill-down"]:
    para(tf, h, 17, NAVY, bold=True, font=HF, first=(h == "Overview"), space_after=3, bullet=True)
para(tf, "", 6, GREY)
para(tf, "One self-contained HTML file. Opens in any browser, fully offline — no install, "
        "no server, no internet. Every decision shows image, ground-truth, score vs. threshold, ✓/✗.",
        15, DARK)
img_fit(s, first_existing(os.path.join(ROOT, "07_paper_draft", "dashboard_snapshots", "dashboard_explainability.png"),
        os.path.join(QUAL, "qualitative_success_cases.png"),
        os.path.join(FIG, "figure_5_failure_cases.png")), 6.4, 1.55, 6.35, 5.1)

# ---------------- S9 Honest positioning ----------------
s = slide(WHITE)
title(s, "Honest about where we stand")
cols = [("Ours — DINOv2 (now)", "0.76", "frozen DINOv2-small,\nfully reproducible", NAVY),
        ("Ours + padding/base", "~0.83–0.90", "mask padding, DINOv2-base,\nmulti-seed", TEAL),
        ("SOTA (CSAD/SALAD)", "0.95–0.96", "component segmentation,\nheavy training", GREY)]
cx = 0.7; cw = 3.95; gap = 0.18
for i, (h, num, d, col) in enumerate(cols):
    l = cx + i * (cw + gap)
    rect(s, l, 1.6, cw, 3.0, PANEL)
    rect(s, l, 1.6, cw, 0.7, col)
    _, tf = box(s, l + 0.15, 1.68, cw - 0.3, 0.55)
    para(tf, h, 16, WHITE, bold=True, font=HF, first=True, align=PP_ALIGN.CENTER, space_after=0)
    _, tf = box(s, l + 0.15, 2.5, cw - 0.3, 2.0)
    para(tf, num, 34, col, bold=True, font=HF, first=True, align=PP_ALIGN.CENTER, space_after=6)
    for j, line in enumerate(d.split("\n")):
        para(tf, line, 14, DARK, align=PP_ALIGN.CENTER, space_after=0, first=False)
_, tf = box(s, 0.7, 4.95, 11.9, 1.9)
para(tf, "Our contribution is rigor, reproducibility, and explainability — not chasing a leaderboard number.",
     18, NAVY, bold=True, font=HF, first=True, space_after=8)
para(tf, "We name each method for what it actually computes, keep evaluation leakage-free, and audit the "
        "data end to end. Features are verified as real DINOv2; the remaining gap is padding-masking, a "
        "larger backbone, and multi-seed — not a redesign.", 15, DARK)

# ---------------- S10 Conclusion ----------------
s = slide(DEEP)
rect(s, 0.0, 0.0, 0.28, SH, TEAL, rounded=False)
title(s, "Summary & next steps", color=WHITE)
_, tf = box(s, 0.9, 1.7, 11.5, 4.8)
para(tf, "A complete, working product: reproducible pipeline, four detectors + fusion, and an "
        "offline explainable dashboard.", 18, WHITE, first=True, space_after=12, bullet=True)
para(tf, "Best fusion: 0.76 overall AUROC (0.73 logical / 0.80 structural) on real DINOv2 features — "
        "a transparent, defensible result.", 18, WHITE, space_after=12, bullet=True)
para(tf, "Feature selection (drop redundant + harmful branches) gives a leaner, better model (0.764).",
     18, WHITE, space_after=12, bullet=True)
para(tf, "Next: mask letterbox padding, move to DINOv2-base, add the official sPRO metric, and report "
        "multi-seed confidence intervals.", 18, TEAL, bold=True, space_after=0, bullet=True)
_, tf = box(s, 0.9, 6.5, 11.5, 0.7)
para(tf, "Thank you — live demo: open dashboard.html", 16, ICE, italic=True, first=True)

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
