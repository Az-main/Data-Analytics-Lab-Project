# Build LOCO Anomaly Inspector — technical presentation v2 (22 slides, dark theme).
# All numbers come from 06_method_results/Final_Evaluation + 02_audit + 04_probe_results.
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------- design system ----------
BG     = RGBColor(0x0E, 0x16, 0x24)
CARD   = RGBColor(0x18, 0x23, 0x3A)
CARD2  = RGBColor(0x1F, 0x2D, 0x4A)
BORDER = RGBColor(0x2A, 0x3A, 0x5C)
TXT    = RGBColor(0xF2, 0xF5, 0xFA)
MUT    = RGBColor(0x9F, 0xB0, 0xC9)
AMBER  = RGBColor(0xF5, 0xA6, 0x23)
BLUE   = RGBColor(0x7F, 0xB3, 0xE8)
GREEN  = RGBColor(0x3D, 0xCB, 0x8F)
RED    = RGBColor(0xE0, 0x56, 0x6B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
HEAD_F = "Trebuchet MS"
BODY_F = "Calibri"

SW, SH = 13.333, 7.5
A = "slide_assets"
F = "figures"

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background(); r.shadow.inherit = False
    return s

def _set_font(run, font=BODY_F, size=13, color=TXT, bold=False, italic=False):
    run.font.name = font; run.font.size = Pt(size)
    run.font.color.rgb = color; run.font.bold = bold; run.font.italic = italic

def _baseline(run, pct):
    """pct: positive = superscript, negative = subscript (e.g. 30 / -25)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set('baseline', str(int(pct * 1000)))

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph = list of run dicts."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line
        for rd in para:
            r = p.add_run()
            r.text = rd.get("t", "")
            _set_font(r, rd.get("f", BODY_F), rd.get("s", 13), rd.get("c", TXT),
                      rd.get("b", False), rd.get("i", False))
            if "base" in rd: _baseline(r, rd["base"])
    return tb

def kicker(s, txt, x=0.62, y=0.40, color=AMBER):
    text(s, x, y, 11.0, 0.3, [[{"t": txt.upper(), "f": HEAD_F, "s": 11, "c": color, "b": True}]])

def title(s, txt, x=0.6, y=0.68, size=29):
    text(s, x, y, 12.1, 0.75, [[{"t": txt, "f": HEAD_F, "s": size, "c": TXT, "b": True}]])

def page(s, n):
    text(s, 12.45, 7.08, 0.7, 0.3, [[{"t": f"{n:02d}", "f": HEAD_F, "s": 10, "c": MUT}]],
         align=PP_ALIGN.RIGHT)

def card(s, x, y, w, h, fill=CARD, border=BORDER):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.adjustments[0] = min(0.07, 0.07 * 2.0 / h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = border; c.line.width = Pt(1)
    c.shadow.inherit = False
    return c

def chip(s, x, y, w, txt, color=AMBER, h=0.34, size=11):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    c.adjustments[0] = 0.5
    c.fill.solid(); c.fill.fore_color.rgb = CARD2
    c.line.color.rgb = color; c.line.width = Pt(1); c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.06); tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    _set_font(r, HEAD_F, size, color, bold=True)
    return c

def pic(s, path, x, y, w=None, h=None):
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(x), Inches(y), **kw)

def arrow(s, x, y, w, h=0.30):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.adjustments[0] = 0.55; a.adjustments[1] = 0.55
    a.fill.solid(); a.fill.fore_color.rgb = BORDER
    a.line.fill.background(); a.shadow.inherit = False
    return a

def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt

def stat(s, x, y, w, big, label, color=AMBER, big_size=30, label_size=11.5):
    text(s, x, y, w, 0.55, [[{"t": big, "f": HEAD_F, "s": big_size, "c": color, "b": True}]])
    text(s, x, y + 0.58, w, 0.6, [[{"t": label, "f": BODY_F, "s": label_size, "c": MUT}]], line=0.95)

# =========================================================================
# SLIDE 1 — Title
# =========================================================================
s = slide()
text(s, 0.9, 1.30, 7.6, 0.3, [[{"t": "DATA ANALYTICS LAB · 8TH TRIMESTER · PROJECT A",
                                "f": HEAD_F, "s": 11, "c": AMBER, "b": True}]])
text(s, 0.9, 1.62, 7.7, 1.0, [[{"t": "LOCO Anomaly Inspector", "f": HEAD_F, "s": 44, "c": TXT, "b": True}]])
text(s, 0.9, 2.62, 7.5, 1.5, [[
    {"t": "Logical & structural anomaly detection on the MVTec LOCO AD industrial benchmark — "
          "a leakage-free DINOv2 pipeline with validated fusion and an explainable, fully offline dashboard.",
     "s": 15.5, "c": MUT}]], line=1.15)
text(s, 0.9, 4.02, 7.5, 0.4, [[
    {"t": "Azmain Islam", "f": HEAD_F, "s": 14, "c": TXT, "b": True},
    {"t": "    ·    ID 0152330058", "s": 12.5, "c": MUT}]])
text(s, 0.9, 4.42, 7.5, 0.4, [[
    {"t": "Shamanta Rahman Roza", "f": HEAD_F, "s": 14, "c": TXT, "b": True},
    {"t": "    ·    ID 0152330045", "s": 12.5, "c": MUT}]])
text(s, 0.9, 4.86, 7.5, 0.35, [[
    {"t": "Data Analytics Lab  ·  June 2026", "s": 11.5, "c": MUT}]])
stat(s, 0.9, 5.35, 2.4, "3,651", "images · 5 categories", big_size=28)
stat(s, 3.35, 5.35, 2.4, "0.761", "fused image AUROC", big_size=28)
stat(s, 5.8, 5.35, 2.7, "100%", "offline, reproducible", big_size=28)
for (img, px, py) in [("ex_breakfast_logical.png", 8.75, 1.45), ("ex_juice_structural.png", 10.95, 1.45),
                      ("ex_pushpins_logical.png", 8.75, 3.65), ("ex_splicing_logical.png", 10.95, 3.65)]:
    pic(s, f"{A}/{img}", px, py, w=2.05, h=2.05)
text(s, 8.75, 5.85, 4.25, 0.5, [[
    {"t": "MVTec LOCO AD test samples — ground-truth defect regions in red", "s": 10.5, "c": MUT, "i": True}]])
notes(s, "Frame the project in one sentence: an end-to-end, auditable anomaly-detection system for the "
         "MVTec LOCO benchmark. Mention the three headline numbers and that everything shown is reproducible "
         "from the repo. 30–40 seconds, then move on.")

# =========================================================================
# SLIDE 2 — Agenda
# =========================================================================
s = slide()
kicker(s, "Agenda")
title(s, "What this talk covers")
agenda = [
    ("01", "Problem & benchmark", "Why logical anomalies break classic AD; the MVTec LOCO dataset."),
    ("02", "Data audit & EDA", "Leakage checks, letterbox preprocessing, and the EDA that shaped the model."),
    ("03", "Methods", "Four frozen-DINOv2 detectors and a validation-only fusion."),
    ("04", "Results & ablations", "Main table, per-category analysis, feature selection, efficiency."),
    ("05", "Explainability & dashboard", "Every decision inspectable in a self-contained offline dashboard."),
    ("06", "Limits & roadmap", "Honest positioning vs SOTA and what comes next."),
]
for i, (num, head, body) in enumerate(agenda):
    cx = 0.6 + (i % 3) * 4.20
    cy = 1.75 + (i // 3) * 2.45
    card(s, cx, cy, 3.95, 2.2)
    text(s, cx + 0.3, cy + 0.25, 3.35, 0.5, [[{"t": num, "f": HEAD_F, "s": 24, "c": AMBER, "b": True}]])
    text(s, cx + 0.3, cy + 0.82, 3.35, 0.4, [[{"t": head, "f": HEAD_F, "s": 15.5, "c": TXT, "b": True}]])
    text(s, cx + 0.3, cy + 1.24, 3.35, 0.85, [[{"t": body, "s": 11.5, "c": MUT}]], line=1.05)
page(s, 2)
notes(s, "One breath per box. Signal that results AND their honest limits both get airtime — that is the "
         "story of the project: rigor over leaderboard-chasing.")

# =========================================================================
# SLIDE 3 — Motivation
# =========================================================================
s = slide()
kicker(s, "Motivation")
title(s, "Why industrial anomaly detection is hard")
points = [
    ("1", "No labels for what can go wrong",
     "Production lines collect thousands of good images but few, unpredictable defects. "
     "Models must train on normal data only — unsupervised anomaly detection."),
    ("2", "Two very different failure modes",
     "Local damage (a scratch, a contamination) and global rule violations (missing, extra or "
     "misplaced parts) leave completely different statistical footprints."),
    ("3", "Decisions must be explainable",
     "A bare score convinces nobody on a factory floor. Operators need the score, the threshold "
     "it is compared against, and the evidence — for every single image."),
]
for i, (num, head, body) in enumerate(points):
    py = 1.75 + i * 1.72
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.62), Inches(py), Inches(0.52), Inches(0.52))
    c.fill.solid(); c.fill.fore_color.rgb = CARD2; c.line.color.rgb = AMBER
    c.line.width = Pt(1.25); c.shadow.inherit = False
    tf = c.text_frame; tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = num; _set_font(r, HEAD_F, 16, AMBER, bold=True)
    text(s, 1.4, py - 0.04, 6.1, 0.4, [[{"t": head, "f": HEAD_F, "s": 16, "c": TXT, "b": True}]])
    text(s, 1.4, py + 0.40, 6.1, 1.1, [[{"t": body, "s": 12.5, "c": MUT}]], line=1.1)
card(s, 7.95, 1.7, 4.78, 5.35)
pic(s, f"{A}/ex_breakfast_logical.png", 8.33, 2.0, w=4.0, h=4.0)
text(s, 8.33, 6.1, 4.0, 0.85, [[
    {"t": "Nothing here is damaged. ", "s": 12, "c": TXT, "b": True},
    {"t": "The compartment contents violate the product's composition rules — a logical anomaly "
          "(ground-truth region in red).", "s": 12, "c": MUT}]], line=1.08)
page(s, 3)
notes(s, "Sell the problem before the solution. Emphasise point 2: it predicts the whole method design — "
         "we will need local AND global detectors. Point at the breakfast box: every patch is normal, "
         "the composition is not.")

# =========================================================================
# SLIDE 4 — Two kinds of anomaly
# =========================================================================
s = slide()
kicker(s, "Problem definition")
title(s, "Two kinds of anomaly, one detector budget")
# Structural card
card(s, 0.6, 1.62, 6.0, 4.42)
text(s, 0.92, 1.86, 5.4, 0.4, [[{"t": "STRUCTURAL", "f": HEAD_F, "s": 15, "c": GREEN, "b": True},
                                {"t": "  — appearance breaks", "f": HEAD_F, "s": 15, "c": TXT, "b": True}]])
pic(s, f"{A}/ex_juice_structural.png", 0.92, 2.36, w=2.55, h=2.55)
text(s, 3.65, 2.36, 2.7, 2.6, [[
    {"t": "Local defects: scratches, dents, contaminations.", "s": 12.5, "c": TXT}], [
    {"t": "Evidence is confined to a small region — comparing patches against normal "
          "appearance finds it.", "s": 12.5, "c": MUT}]], line=1.1, space_after=8)
text(s, 0.92, 5.05, 5.4, 0.8, [[
    {"t": "juice_bottle — a local defect near the bottle neck (red overlay). A patch-level "
          "nearest-neighbour search nails this.", "s": 11, "c": MUT, "i": True}]], line=1.05)
# Logical card
card(s, 6.85, 1.62, 6.0, 4.42)
text(s, 7.17, 1.86, 5.4, 0.4, [[{"t": "LOGICAL", "f": HEAD_F, "s": 15, "c": AMBER, "b": True},
                                {"t": "  — rules break", "f": HEAD_F, "s": 15, "c": TXT, "b": True}]])
pic(s, f"{A}/ex_pushpins_logical.png", 7.17, 2.36, w=2.55, h=2.55)
text(s, 9.9, 2.36, 2.7, 2.6, [[
    {"t": "Wrong count, position or combination of otherwise-normal parts.", "s": 12.5, "c": TXT}], [
    {"t": "Every patch looks perfectly normal in isolation; only the global composition "
          "is wrong.", "s": 12.5, "c": MUT}]], line=1.1, space_after=8)
text(s, 7.17, 5.05, 5.4, 0.8, [[
    {"t": "pushpins — extra / misplaced pins highlighted in red. No patch is locally "
          "abnormal.", "s": 11, "c": MUT, "i": True}]], line=1.05)
card(s, 0.6, 6.28, 12.25, 0.78, fill=CARD2)
text(s, 0.92, 6.45, 11.7, 0.5, [[
    {"t": "Why LOCO exists:  ", "f": HEAD_F, "s": 12.5, "c": AMBER, "b": True},
    {"t": "patch-based detectors (PatchCore-style) excel at structural defects but are nearly blind to "
          "logical ones (Bergmann et al., IJCV 2022). Closing that gap drives every design choice that follows.",
     "s": 12.5, "c": TXT}]])
page(s, 4)
notes(s, "The core conceptual slide. Make the audience predict the method: local memory for structural, "
         "something global for logical. If they get this, the rest of the talk is easy.")

# =========================================================================
# SLIDE 5 — Dataset
# =========================================================================
s = slide()
kicker(s, "Benchmark")
title(s, "MVTec LOCO AD — five products, three splits")
stat(s, 0.6, 1.62, 1.9, "3,651", "images total", big_size=26)
stat(s, 2.65, 1.62, 1.9, "1,778", "train — good only", big_size=26)
stat(s, 4.7, 1.62, 1.9, "305", "validation — good only", big_size=26)
stat(s, 6.75, 1.62, 1.9, "1,568", "held-out test", big_size=26)
rows = [
    ("Category", "Train", "Val", "Test", "Test composition"),
    ("breakfast_box", "351", "62", "275", "102 good · 83 logical · 90 structural"),
    ("juice_bottle", "335", "54", "330", "94 good · 142 logical · 94 structural"),
    ("pushpins", "372", "69", "310", "138 good · 91 logical · 81 structural"),
    ("screw_bag", "360", "60", "341", "122 good · 137 logical · 82 structural"),
    ("splicing_connectors", "360", "60", "312", "119 good · 108 logical · 85 structural"),
]
tx, ty, tw, th = 0.6, 3.0, 8.05, 3.3
tbl_shape = s.shapes.add_table(len(rows), 5, Inches(tx), Inches(ty), Inches(tw), Inches(th))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(0.85)
tbl.columns[2].width = Inches(0.85)
tbl.columns[3].width = Inches(0.85)
tbl.columns[4].width = Inches(3.5)
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD2 if ri == 0 else (CARD if ri % 2 else RGBColor(0x14, 0x1E, 0x32))
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.05)
        cell.margin_top = cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = val
        _set_font(r, BODY_F if ri else HEAD_F, 11.5 if ri else 12,
                  TXT if ri == 0 else (TXT if ci == 0 else MUT), bold=(ri == 0))
        if 1 <= ci <= 3: p.alignment = PP_ALIGN.CENTER
text(s, 0.6, 6.5, 8.0, 0.5, [[
    {"t": "Validation contains only good images — anomaly types are unknown until test time, by design.",
     "s": 11.5, "c": MUT, "i": True}]])
for (img, py, cap) in [("ex_screwbag_logical.png", 1.62, "screw_bag — wrong part set (red)"),
                       ("ex_splicing_logical.png", 4.32, "splicing_connectors — wiring rule violation (red)")]:
    pic(s, f"{A}/{img}", 9.05, py, w=2.45, h=2.45)
    text(s, 11.6, py + 0.7, 1.5, 1.4, [[{"t": cap, "s": 11, "c": MUT, "i": True}]], line=1.05)
page(s, 5)
notes(s, "Key protocol point: the validation split has NO anomalies. Every threshold and every fusion weight "
         "must therefore come from good images only — this constraint shapes the whole calibration story.")

# =========================================================================
# SLIDE 6 — Data audit & preprocessing
# =========================================================================
s = slide()
kicker(s, "Rigor first")
title(s, "Data audit & preprocessing — before any modeling")
steps = [
    ("01", "Raw data is read-only",
     "Full inventory of 3,651 product images + 1,246 defect masks; per-file MD5 manifest stored as CSV. "
     "The raw tree is never modified."),
    ("02", "Aspect-preserving letterbox to 384×384",
     "Bilinear for images, nearest-neighbour for masks (labels must stay binary). Scale and padding "
     "offsets logged per image — every pixel is traceable back."),
    ("03", "Cross-split leakage check",
     "MD5 + perceptual hashing compared across train / validation / test. "
     "Result: zero duplicate leaks. The headline numbers cannot be inflated by memorised images."),
    ("04", "Split discipline, enforced",
     "Fit on Train(good). Calibrate on Validation(good). Test stays untouched until the final, single "
     "evaluation pass. No test labels in any tuning loop."),
]
for i, (num, head, body) in enumerate(steps):
    cx = 0.6 + (i % 2) * 6.25
    cy = 1.72 + (i // 2) * 2.32
    card(s, cx, cy, 6.0, 2.1)
    text(s, cx + 0.3, cy + 0.22, 0.9, 0.5, [[{"t": num, "f": HEAD_F, "s": 22, "c": AMBER, "b": True}]])
    text(s, cx + 1.1, cy + 0.24, 4.6, 0.4, [[{"t": head, "f": HEAD_F, "s": 14.5, "c": TXT, "b": True}]])
    text(s, cx + 1.1, cy + 0.70, 4.6, 1.3, [[{"t": body, "s": 11.8, "c": MUT}]], line=1.08)
text(s, 0.6, 6.45, 12.2, 0.5, [[
    {"t": "Every step emits CSV artifacts (hashes, counts, resize metadata) — the audit is reproducible, "
          "not anecdotal.", "s": 12.5, "c": AMBER, "i": True}]])
page(s, 6)
notes(s, "This is the slide reviewers tend to reward: most course projects skip the audit. Stress the "
         "zero-leakage finding and that masks use nearest-neighbour interpolation — a classic silent bug "
         "elsewhere.")

# =========================================================================
# SLIDE 7 — EDA
# =========================================================================
s = slide()
kicker(s, "Exploratory analysis")
title(s, "EDA that informed the model — not decoration")
findings = [
    ("Defect footprint separates the two types",
     "Structural masks are small and local (median ≈ 1–2% of the image). Logical masks can cover "
     "over 30% (screw_bag). One detector cannot serve both — we need local and global branches."),
    ("Anomalies live on the product, not the padding",
     "Defect-frequency heatmaps concentrate on the object region. Letterbox borders carry no defect "
     "signal — yet still feed the feature extractor. We track this as a known caveat."),
    ("Five very different layouts",
     "Fixed grid (pushpins) vs free-form tray (breakfast_box): no single spatial prior fits all. "
     "Position must be modeled explicitly — this motivates the region-aware memory."),
]
for i, (head, body) in enumerate(findings):
    py = 1.70 + i * 1.62
    text(s, 0.6, py, 7.0, 0.4, [[{"t": head, "f": HEAD_F, "s": 14.5, "c": BLUE, "b": True}]])
    text(s, 0.6, py + 0.42, 7.0, 1.1, [[{"t": body, "s": 12.2, "c": MUT}]], line=1.08)
text(s, 0.6, 6.62, 7.0, 0.5, [[
    {"t": "Test anomalies are used for descriptive EDA only — never for model tuning.",
     "s": 11.5, "c": AMBER, "i": True}]])
card(s, 7.95, 1.70, 4.78, 5.0)
pic(s, f"{A.replace('slide_assets','../04_probe_results') if False else '../04_probe_results'}/eda_spatial_heatmap_cleaned.png",
    8.23, 1.95, w=4.22, h=3.52)
text(s, 8.23, 5.55, 4.2, 1.0, [[
    {"t": "Where defects occur (all categories pooled): ", "s": 11.5, "c": TXT, "b": True},
    {"t": "frequency mass sits on the product region; padding stays dark.", "s": 11.5, "c": MUT}]], line=1.08)
page(s, 7)
notes(s, "Each finding maps 1:1 to a design decision — say that mapping out loud. Finding 1 → two branch "
         "families; finding 2 → padding caveat (returns on the limitations slide); finding 3 → region-aware "
         "memory.")

# =========================================================================
# SLIDE 8 — Negative result probe
# =========================================================================
s = slide()
kicker(s, "A negative result we kept")
title(s, "Probe: can we just segment & count components? No.")
text(s, 0.6, 1.70, 7.4, 0.7, [[
    {"t": "Hypothesis: ", "f": HEAD_F, "s": 14, "c": BLUE, "b": True},
    {"t": "if components could be segmented and counted reliably, logical anomalies would reduce to "
          "comparing counts.", "s": 13, "c": TXT}]], line=1.1)
probe = [
    ("Method", "PCA-RGB projection + K-means clustering on frozen DINOv2 patch features; connected "
               "components treated as candidate 'parts'."),
    ("Result", "Counts are unstable even between two normal images of the same product — 6 vs 1 components "
               "on near-identical breakfast trays. Clusters track texture, not objects."),
    ("Decision", "Reject explicit component counting as the main method. Encode composition statistically "
                 "instead: visual-word histograms + region-aware memories."),
]
for i, (head, body) in enumerate(probe):
    py = 2.55 + i * 1.30
    chip(s, 0.6, py, 1.25, head.upper(), color=AMBER if head == "Decision" else BLUE)
    text(s, 2.05, py - 0.02, 5.95, 1.2, [[{"t": body, "s": 12.2, "c": MUT}]], line=1.08)
text(s, 0.6, 6.55, 7.4, 0.6, [[
    {"t": "Why show a failure? ", "f": HEAD_F, "s": 12.5, "c": AMBER, "b": True},
    {"t": "Cheap probes killed a dead-end before it became the project. Negative results steer design.",
     "s": 12.5, "c": TXT}]])
card(s, 8.35, 1.55, 4.38, 5.5, fill=WHITE, border=BORDER)
pic(s, f"{A}/probe_crop.png", 8.55, 1.75, w=3.98, h=4.60)
text(s, 8.55, 6.40, 3.98, 0.6, [[
    {"t": "image · PCA-RGB · components — two normal trays yield counts 6 vs 1",
     "s": 10.5, "c": RGBColor(0x44, 0x4C, 0x5A), "i": True}]], line=1.0)
page(s, 8)
notes(s, "This slide demonstrates scientific process. One sentence of humility ('our first idea failed') "
         "buys a lot of credibility for the numbers that follow.")

# =========================================================================
# SLIDE 9 — Pipeline overview
# =========================================================================
s = slide()
kicker(s, "Method overview")
title(s, "One pipeline, end to end — every stage auditable")
stages = [
    ("Audit + letterbox", "read-only raw,\nMD5 manifest,\n384×384 letterbox"),
    ("Frozen DINOv2-small", "patch tokens,\nd = 384 —\nzero training"),
    ("Four anomaly scorers", "2 local memories,\n2 global\ndescriptors"),
    ("Calibrate & fuse", "z-norm + mean rule,\nValidation(good)\nonly"),
    ("Evaluate & explain", "held-out test,\nAUROC / F1 →\ndashboard"),
]
bw, gap = 2.28, 0.24
x0 = (SW - (bw * 5 + gap * 4)) / 2
for i, (head, body) in enumerate(stages):
    bx = x0 + i * (bw + gap)
    card(s, bx, 1.80, bw, 1.95, fill=CARD2 if i in (1, 3) else CARD)
    text(s, bx + 0.16, 1.98, bw - 0.32, 0.65,
         [[{"t": head, "f": HEAD_F, "s": 13.5, "c": AMBER if i in (1, 3) else TXT, "b": True}]], line=0.98)
    text(s, bx + 0.16, 2.62, bw - 0.32, 1.0,
         [[{"t": body.replace("\n", " "), "s": 11, "c": MUT}]], line=1.05)
    if i < 4:
        arrow(s, bx + bw + 0.015, 2.62, gap - 0.03, 0.26)
dets = [
    ("Patch Memory", "local appearance", BLUE),
    ("Region-Aware Memory", "local + position", AMBER),
    ("Composition Histogram", "global counts", GREEN),
    ("Image-Stat (proxy)", "global baseline", MUT),
]
text(s, x0, 4.05, 11.0, 0.35, [[{"t": "THE FOUR SCORERS", "f": HEAD_F, "s": 11, "c": MUT, "b": True}]])
for i, (head, sub, col) in enumerate(dets):
    bx = 0.6 + i * 3.075
    card(s, bx, 4.42, 2.92, 1.05)
    text(s, bx + 0.2, 4.55, 2.6, 0.35, [[{"t": head, "f": HEAD_F, "s": 12.5, "c": col, "b": True}]])
    text(s, bx + 0.2, 4.93, 2.6, 0.35, [[{"t": sub, "s": 11, "c": MUT}]])
card(s, 0.6, 5.85, 12.25, 1.05, fill=CARD2)
text(s, 0.92, 6.04, 11.7, 0.75, [[
    {"t": "Three rules everywhere:   ", "f": HEAD_F, "s": 13, "c": AMBER, "b": True},
    {"t": "backbone stays frozen (zero training)  ·  test data never enters fitting or tuning  ·  "
          "every stage writes CSV artifacts that the report and dashboard are rebuilt from.",
     "s": 13, "c": TXT}]], line=1.1)
page(s, 9)
notes(s, "Walk left to right in ~45 seconds. The two highlighted boxes are where this project differs from "
         "a baseline run: frozen-feature scorers and validation-only calibration.")

# =========================================================================
# SLIDE 10 — Detectors 1 & 2
# =========================================================================
s = slide()
kicker(s, "Methods · local branches")
title(s, "Detectors 1–2: nearest-neighbour patch memories")
card(s, 0.6, 1.62, 6.0, 4.55)
text(s, 0.92, 1.84, 5.4, 0.4, [[{"t": "D1 · DINOv2 Patch Memory", "f": HEAD_F, "s": 15.5, "c": BLUE, "b": True},
                                {"t": "   (PatchCore-style)", "s": 12, "c": MUT}]])
text(s, 0.92, 2.34, 5.4, 1.9, [
    [{"t": "Build: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "memory bank M of normal patch tokens (coreset, 50,000 entries, d = 384).", "s": 12.3, "c": MUT}],
    [{"t": "Score: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "each test patch → cosine distance to its nearest normal patch; the image score is the worst "
           "patch.", "s": 12.3, "c": MUT}],
    [{"t": "Catches structural: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "a damaged patch has no close normal neighbour.", "s": 12.3, "c": MUT}],
], line=1.08, space_after=7)
card(s, 0.92, 4.62, 5.35, 0.62, fill=BG, border=BORDER)
text(s, 1.12, 4.74, 5.0, 0.4, [[
    {"t": "s(x) = max", "f": BODY_F, "s": 13, "c": GREEN},
    {"t": "p", "f": BODY_F, "s": 13, "c": GREEN, "base": -25},
    {"t": " min", "f": BODY_F, "s": 13, "c": GREEN},
    {"t": "m∈M", "f": BODY_F, "s": 13, "c": GREEN, "base": -25},
    {"t": " [ 1 − cos(f", "f": BODY_F, "s": 13, "c": GREEN},
    {"t": "p", "f": BODY_F, "s": 13, "c": GREEN, "base": -25},
    {"t": ", m) ]", "f": BODY_F, "s": 13, "c": GREEN}]])
chip(s, 0.92, 5.55, 1.7, "AUROC 0.724", color=BLUE)
chip(s, 2.78, 5.55, 1.85, "367 ms / image", color=MUT)
card(s, 6.85, 1.62, 6.0, 4.55)
text(s, 7.17, 1.84, 5.4, 0.4, [[{"t": "D2 · Region-Aware Memory", "f": HEAD_F, "s": 15.5, "c": AMBER, "b": True},
                                {"t": "   (strongest single)", "s": 12, "c": MUT}]])
text(s, 7.17, 2.34, 5.4, 2.1, [
    [{"t": "Build: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "partition the token grid into 6×6 spatial regions; one memory per region "
           "(36 per category, 455k entries total).", "s": 12.3, "c": MUT}],
    [{"t": "Score: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "a patch is compared only to normal patches from the same region — appearance and position "
           "must both match; mean of top-5 neighbour distances.", "s": 12.3, "c": MUT}],
    [{"t": "Catches: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "“right part, wrong place” — the positional half of logical anomalies.", "s": 12.3, "c": MUT}],
], line=1.08, space_after=7)
chip(s, 7.17, 5.55, 1.7, "AUROC 0.750", color=AMBER)
chip(s, 9.03, 5.55, 1.85, "154 ms / image", color=MUT)
text(s, 0.6, 6.50, 12.2, 0.5, [[
    {"t": "All hyperparameters (region size 6, top-5) selected on validation-good stability — "
          "details on the ablation slide.", "s": 12, "c": MUT, "i": True}]])
page(s, 10)
notes(s, "D1 is the literature workhorse — give it 30 seconds. Spend the time on D2: same memory idea, but "
         "conditioned on position, which is exactly what misplaced-part anomalies need. It is also 2.4x "
         "faster because each query searches a 36x smaller memory.")

# =========================================================================
# SLIDE 11 — Detectors 3 & 4
# =========================================================================
s = slide()
kicker(s, "Methods · global branches")
title(s, "Detectors 3–4: composition without segmentation")
card(s, 0.6, 1.62, 6.0, 4.40)
text(s, 0.92, 1.84, 5.4, 0.4, [[{"t": "D3 · Composition Histogram", "f": HEAD_F, "s": 15.5, "c": GREEN, "b": True}]])
text(s, 0.92, 2.34, 5.4, 1.85, [
    [{"t": "Vocabulary: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "k-means over normal patch tokens → 32 visual words.", "s": 12.3, "c": MUT}],
    [{"t": "Represent: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "each image → normalised word histogram h(x): what is present, in what proportion.", "s": 12.3, "c": MUT}],
    [{"t": "Score: ", "s": 12.3, "c": TXT, "b": True},
     {"t": "Mahalanobis distance to the normal histogram distribution — counting without segmenting.",
      "s": 12.3, "c": MUT}],
], line=1.08, space_after=7)
card(s, 0.92, 4.42, 5.35, 0.62, fill=BG, border=BORDER)
text(s, 1.12, 4.54, 5.0, 0.4, [[
    {"t": "s(x) = √( (h−μ)", "f": BODY_F, "s": 13, "c": GREEN},
    {"t": "T", "f": BODY_F, "s": 13, "c": GREEN, "base": 30},
    {"t": " Σ", "f": BODY_F, "s": 13, "c": GREEN},
    {"t": "−1", "f": BODY_F, "s": 13, "c": GREEN, "base": 30},
    {"t": " (h−μ) )", "f": BODY_F, "s": 13, "c": GREEN}]])
chip(s, 0.92, 5.38, 1.95, "AUROC 0.576 alone", color=RED)
chip(s, 3.03, 5.38, 1.8, "28 ms / image", color=MUT)
card(s, 6.85, 1.62, 6.0, 4.40)
text(s, 7.17, 1.84, 5.4, 0.4, [[{"t": "D4 · Global Image-Stat", "f": HEAD_F, "s": 15.5, "c": MUT, "b": True},
                                {"t": "   (honestly-named proxy)", "s": 12, "c": MUT}]])
text(s, 7.17, 2.34, 5.4, 1.85, [
    [{"t": "One global descriptor per image, scored against the normal mean — an EfficientAD-style "
           "stand-in, deliberately labeled a proxy rather than passed off as the real network.",
      "s": 12.3, "c": MUT}],
    [{"t": "The only non-DINOv2 branch, and the fastest.", "s": 12.3, "c": MUT}],
    [{"t": "Weak alone — but the most complementary branch in fusion: removing it costs −0.043 AUROC "
           "(leave-one-out).", "s": 12.3, "c": TXT, "b": True}],
], line=1.08, space_after=7)
chip(s, 7.17, 5.38, 1.7, "AUROC 0.655", color=MUT)
chip(s, 9.03, 5.38, 1.75, "21 ms / image", color=MUT)
card(s, 0.6, 6.30, 12.25, 0.76, fill=CARD2)
text(s, 0.92, 6.46, 11.7, 0.5, [[
    {"t": "Design logic:  ", "f": HEAD_F, "s": 12.5, "c": AMBER, "b": True},
    {"t": "two local + two global branches that fail differently. Fusion exploits their disagreement — "
          "single-branch weakness is not a verdict.", "s": 12.5, "c": TXT}]])
page(s, 11)
notes(s, "Two honesty beats live here: D4 is explicitly a proxy (named as such everywhere), and D3's weak "
         "standalone number is shown in red — it sets up the feature-selection slide where it turns out to "
         "actively hurt the fusion.")

# =========================================================================
# SLIDE 12 — Fusion & protocol
# =========================================================================
s = slide()
kicker(s, "Methods · calibration")
title(s, "Leakage-free fusion: calibrated on good images only")
steps12 = [
    ("1", "z-normalise each branch on Validation(good)",
     "Branch scores live on different scales; statistics (μ, σ) come from good validation images only."),
    ("2", "Combine with the pre-specified mean rule",
     "The mean rule was fixed before test evaluation. Max, rank-average and weighted variants are "
     "reported as ablations — weighted ties mean at 0.761."),
    ("3", "Threshold at the 90th percentile of fused val-good scores",
     "This single operating point drives every ✓/✗ verdict in the dashboard and the F1 numbers."),
]
for i, (num, head, body) in enumerate(steps12):
    py = 1.72 + i * 1.38
    text(s, 0.6, py, 0.5, 0.5, [[{"t": num, "f": HEAD_F, "s": 22, "c": AMBER, "b": True}]])
    text(s, 1.25, py + 0.02, 6.0, 0.6, [[{"t": head, "f": HEAD_F, "s": 13.5, "c": TXT, "b": True}]], line=1.0)
    text(s, 1.25, py + 0.60, 6.0, 0.7, [[{"t": body, "s": 11.8, "c": MUT}]], line=1.05)
card(s, 0.6, 5.95, 6.65, 0.72, fill=BG, border=BORDER)
text(s, 0.9, 6.12, 6.2, 0.4, [[
    {"t": "z", "f": BODY_F, "s": 13.5, "c": GREEN},
    {"t": "k", "f": BODY_F, "s": 13.5, "c": GREEN, "base": -25},
    {"t": "(x) = ( s", "f": BODY_F, "s": 13.5, "c": GREEN},
    {"t": "k", "f": BODY_F, "s": 13.5, "c": GREEN, "base": -25},
    {"t": "(x) − μ", "f": BODY_F, "s": 13.5, "c": GREEN},
    {"t": "k", "f": BODY_F, "s": 13.5, "c": GREEN, "base": -25},
    {"t": " ) / σ", "f": BODY_F, "s": 13.5, "c": GREEN},
    {"t": "k", "f": BODY_F, "s": 13.5, "c": GREEN, "base": -25},
    {"t": "      s", "f": BODY_F, "s": 13.5, "c": GREEN},
    {"t": "fused", "f": BODY_F, "s": 13.5, "c": GREEN, "base": -25},
    {"t": "(x) = mean", "f": BODY_F, "s": 13.5, "c": GREEN},
    {"t": "k", "f": BODY_F, "s": 13.5, "c": GREEN, "base": -25},
    {"t": " z", "f": BODY_F, "s": 13.5, "c": GREEN},
    {"t": "k", "f": BODY_F, "s": 13.5, "c": GREEN, "base": -25},
    {"t": "(x)", "f": BODY_F, "s": 13.5, "c": GREEN}]])
card(s, 7.7, 1.72, 5.05, 4.2)
pic(s, f"{A}/chart_fusion_rules.png", 8.0, 1.95, w=4.45, h=2.57)
text(s, 8.0, 4.62, 4.45, 1.1, [[
    {"t": "Fusion-rule ablation. ", "s": 11.8, "c": TXT, "b": True},
    {"t": "Mean (pre-specified) and validation-weighted average tie at 0.761; max and rank-average trail. "
          "No rule was picked using test labels.", "s": 11.8, "c": MUT}]], line=1.08)
chips12 = [("FIT — Train(good)", BLUE), ("CALIBRATE — Val(good)", AMBER), ("REPORT — Test, untouched", GREEN)]
cx = 7.7
for t12, col in chips12:
    wch = 1.62 if "FIT" in t12 else (1.95 if "CALIB" in t12 else 2.18)
    chip(s, cx, 6.30, wch, t12, color=col, h=0.4, size=10)
    cx += wch + 0.22
page(s, 12)
notes(s, "If the audience remembers one methods detail, make it this: nothing downstream of the train/val "
         "split ever sees a test label. The fusion rule itself was committed to in advance — the ablation "
         "exists for completeness, not cherry-picking.")

# =========================================================================
# SLIDE 13 — Main results
# =========================================================================
s = slide()
kicker(s, "Results")
title(s, "Main results — fusion beats every single branch")
card(s, 0.6, 1.62, 8.5, 4.65)
pic(s, f"{A}/chart_main_results.png", 0.85, 1.85, w=8.0, h=3.91)
text(s, 0.85, 5.82, 8.0, 0.4, [[
    {"t": "Image-level AUROC, mean over 5 categories. PatchCore omitted: byte-identical to Patch Memory "
          "in this setup.", "s": 10.5, "c": MUT, "i": True}]])
tk = [
    ("0.761 — fusion wins", AMBER,
     "vs 0.750 for the best single branch. Complementary local + global signals, fused with "
     "validation-only statistics."),
    ("Structural 0.804 / logical 0.731", BLUE,
     "The logical side remains harder — exactly the gap the LOCO benchmark was built to expose."),
    ("Real, verified features", GREEN,
     "Runtime logs record feature_backend = dinov2-small for every branch; deterministic config, "
     "fixed seeds."),
]
for i, (head, col, body) in enumerate(tk):
    py = 1.70 + i * 1.70
    text(s, 9.4, py, 3.35, 0.6, [[{"t": head, "f": HEAD_F, "s": 14, "c": col, "b": True}]], line=1.0)
    text(s, 9.4, py + 0.62, 3.35, 1.0, [[{"t": body, "s": 11.3, "c": MUT}]], line=1.06)
text(s, 0.6, 6.55, 12.2, 0.5, [[
    {"t": "Every number regenerates from 06_method_results/Final_Evaluation/*.csv — the deck, report and "
          "dashboard all read the same tables.", "s": 11.5, "c": MUT, "i": True}]])
page(s, 13)
notes(s, "Pause on this slide. Read the fusion bar triplet once, then the two takeaways on the right. "
         "Anticipate the question 'why is 0.76 good?' — answer comes two slides later in the SOTA "
         "positioning slide.")

# =========================================================================
# SLIDE 14 — Per-category
# =========================================================================
s = slide()
kicker(s, "Results · per category")
title(s, "Where it works, where it struggles — and why")
card(s, 0.6, 1.62, 6.85, 4.35)
pic(s, f"{A}/chart_percat_heatmap.png", 0.92, 1.92, w=6.2, h=3.34)
text(s, 0.92, 5.32, 6.2, 0.5, [[
    {"t": "Best Fusion, image-level AUROC on the held-out test split.", "s": 10.5, "c": MUT, "i": True}]])
tk14 = [
    ("juice_bottle & breakfast_box: strong", GREEN,
     "0.840 / 0.820 overall — rich appearance cues suit patch memories; fill levels and fruit "
     "composition give global signal."),
    ("pushpins: a split personality", AMBER,
     "Structural 0.983 (near-perfect) vs logical 0.600. Detecting a damaged pin ≠ noticing one pin "
     "too many — exact counting is still out of reach."),
    ("screw_bag & splicing: the open problem", RED,
     "Logical 0.628 / 0.676. Small interchangeable parts in variable layouts need component-aware "
     "reasoning — our roadmap's P2."),
]
for i, (head, col, body) in enumerate(tk14):
    py = 1.70 + i * 1.72
    text(s, 7.75, py, 5.0, 0.4, [[{"t": head, "f": HEAD_F, "s": 13.5, "c": col, "b": True}]], line=1.0)
    text(s, 7.75, py + 0.44, 5.0, 1.15, [[{"t": body, "s": 11.6, "c": MUT}]], line=1.07)
text(s, 0.6, 6.55, 12.2, 0.5, [[
    {"t": "Consistent with the literature: across published methods, logical anomalies in counting-heavy "
          "categories remain the hardest cell of this benchmark.", "s": 11.5, "c": MUT, "i": True}]])
page(s, 14)
notes(s, "The pushpins row is the most quotable: 0.983 vs 0.600 in one category is the cleanest possible "
         "demonstration that structural and logical anomalies are different problems.")

# =========================================================================
# SLIDE 15 — Feature selection
# =========================================================================
s = slide()
kicker(s, "Results · ablations")
title(s, "Feature selection: fewer branches, better model")
card(s, 0.6, 1.62, 6.0, 3.55)
text(s, 0.9, 1.80, 5.4, 0.35, [[{"t": "Leave-one-out: what each branch adds",
                                 "f": HEAD_F, "s": 13, "c": TXT, "b": True}]])
pic(s, f"{A}/chart_loo_importance.png", 0.9, 2.25, w=5.4, h=2.36)
card(s, 6.85, 1.62, 6.0, 3.55)
text(s, 7.15, 1.80, 5.4, 0.35, [[{"t": "Branch subsets: overall AUROC",
                                  "f": HEAD_F, "s": 13, "c": TXT, "b": True}]])
pic(s, f"{A}/chart_feature_selection.png", 7.15, 2.30, w=5.4, h=2.55)
card(s, 0.6, 5.45, 12.25, 1.55, fill=CARD2)
text(s, 0.92, 5.62, 11.7, 1.3, [
    [{"t": "Drop the redundant duplicate + the harmful branch → 3-branch model: 0.764 overall, "
           "0.826 structural — better than all five.", "f": HEAD_F, "s": 13.5, "c": AMBER, "b": True}],
    [{"t": "Composition Histogram hurts the fusion (−0.013) despite decent logical-only skill; the Global "
           "proxy is weak alone (0.655) yet most complementary (+0.043). Lesson: single-branch strength "
           "≠ fusion value.", "s": 12.3, "c": TXT}],
], line=1.12, space_after=6)
page(s, 15)
notes(s, "This is the analytical heart of the deck — feature selection done on representations, evaluated "
         "leakage-free. The counter-intuitive pair (weak-but-complementary vs decent-but-harmful) usually "
         "gets the best questions.")

# =========================================================================
# SLIDE 16 — Hyperparameters & efficiency
# =========================================================================
s = slide()
kicker(s, "Results · ablations & cost")
title(s, "Hyperparameters picked blind; cost measured")
card(s, 0.6, 1.62, 5.6, 4.6)
text(s, 0.92, 1.84, 5.0, 0.7, [[
    {"t": "Selected on validation-good stability — never on test", "f": HEAD_F, "s": 14, "c": TXT, "b": True}]],
     line=1.0)
hps = [
    ("Region grid", "4 ·", " 6 ", "· 8"),
    ("Neighbours top-k", "1 ·", " 5 ", "· 10 · mean · max"),
    ("Visual words k", "16 ·", " 32 ", "· 64 · 128"),
    ("Fusion rule", "", " mean ", "— pre-specified, no search"),
]
for i, (label, pre, sel, post) in enumerate(hps):
    py = 2.62 + i * 0.78
    text(s, 0.92, py, 2.1, 0.4, [[{"t": label, "s": 12.5, "c": MUT}]])
    text(s, 3.05, py, 3.0, 0.4, [[
        {"t": pre, "s": 12.5, "c": MUT},
        {"t": sel, "f": HEAD_F, "s": 13, "c": AMBER, "b": True},
        {"t": post, "s": 12.5, "c": MUT}]])
text(s, 0.92, 5.72, 5.0, 0.4, [[
    {"t": "Chosen values highlighted; full grids in the ablation CSVs.", "s": 11, "c": MUT, "i": True}]])
card(s, 6.45, 1.62, 6.4, 4.6)
text(s, 6.77, 1.80, 5.8, 0.35, [[{"t": "Accuracy vs inference cost (CPU)",
                                  "f": HEAD_F, "s": 13, "c": TXT, "b": True}]])
pic(s, f"{A}/chart_speed_accuracy.png", 6.77, 2.22, w=5.75, h=3.10)
text(s, 6.77, 5.40, 5.8, 0.7, [[
    {"t": "Region-Aware sits on the sweet spot: best single AUROC at 2.4× less cost than Patch Memory "
          "(each query searches one region's memory, not all of it).", "s": 11.3, "c": MUT}]], line=1.06)
text(s, 0.6, 6.50, 12.2, 0.5, [[
    {"t": "No GPU training anywhere: frozen features + classical statistics. Full fusion ≈ 0.57 s per image "
          "on CPU.", "s": 12, "c": AMBER, "i": True}]])
page(s, 16)
notes(s, "Two messages: (1) hyperparameter selection cannot have leaked test signal — grids were judged on "
         "validation-good stability; (2) the system is deployable on CPU, which matters for factory edge "
         "hardware.")

# =========================================================================
# SLIDE 17 — Explainability
# =========================================================================
s = slide()
kicker(s, "Explainability")
title(s, "Every decision is inspectable — successes and failures")
pic(s, f"{F}/dashboard_explainability.png", 0.6, 1.62, w=6.05, h=5.28)
tk17 = [
    ("What each card shows", TXT,
     "Input image · ground-truth defect overlay (red) · fused score vs the validation threshold · "
     "verdict checked against the true label (✓/✗)."),
    ("Failures are first-class citizens", AMBER,
     "Example in row 1: score 1.22 vs threshold 1.88 → a missed structural defect, shown in red — "
     "visible, not hidden in an appendix."),
    ("Why it matters", BLUE,
     "Operators can audit any flag in seconds; we used the same view for error analysis "
     "(near-threshold breakfast_box false positives)."),
]
for i, (head, col, body) in enumerate(tk17):
    py = 1.80 + i * 1.66
    text(s, 7.0, py, 5.75, 0.4, [[{"t": head, "f": HEAD_F, "s": 14, "c": col, "b": True}]])
    text(s, 7.0, py + 0.44, 5.75, 1.1, [[{"t": body, "s": 12, "c": MUT}]], line=1.07)
text(s, 7.0, 6.80, 5.75, 0.4, [[
    {"t": "Per-image decisions, Explainability tab of the dashboard.", "s": 10.5, "c": MUT, "i": True}]])
page(s, 17)
notes(s, "Point physically at one card and read it: image, truth, score vs threshold, verdict. Showing a "
         "miss on your own slide is the credibility move of the talk.")

# =========================================================================
# SLIDE 18 — Dashboard
# =========================================================================
s = slide()
kicker(s, "Deliverable")
title(s, "A self-contained, offline analytics dashboard")
pic(s, f"{A}/dashboard_overview_crop.png", 0.6, 1.62, w=7.35, h=4.78)
tk18 = [
    ("One HTML file", "No install, no server, no internet. Opens in any browser — submission-proof."),
    ("Seven tabs", "Overview · EDA · Leaderboard · Feature Selection · Explainability · Per-category · "
                   "Key Insights."),
    ("Transparency banner", "States up front: real frozen DINOv2-small features, single seed, padding "
                            "caveat. Honesty is part of the UI."),
    ("Always in sync", "Rebuilt directly from the result CSVs by script — the dashboard cannot drift from "
                       "the numbers."),
]
for i, (head, body) in enumerate(tk18):
    py = 1.72 + i * 1.18
    text(s, 8.3, py, 4.45, 0.35, [[{"t": head, "f": HEAD_F, "s": 13.5, "c": AMBER if i == 0 else TXT, "b": True}]])
    text(s, 8.3, py + 0.38, 4.45, 0.7, [[{"t": body, "s": 11.3, "c": MUT}]], line=1.05)
text(s, 0.6, 6.60, 7.35, 0.4, [[
    {"t": "Live demo: 09_dashboard/dashboard.html", "f": HEAD_F, "s": 13, "c": GREEN, "b": True}]])
page(s, 18)
notes(s, "If a live demo is allowed, switch to the real dashboard here (it opens offline). Otherwise walk "
         "the screenshot: stat cards, dataset table, methods text — all generated from the same CSVs as "
         "this deck.")

# =========================================================================
# SLIDE 19 — Positioning vs SOTA
# =========================================================================
s = slide()
kicker(s, "Honest positioning")
title(s, "Where 0.761 sits — and what the gap is made of")
cards19 = [
    ("OURS — TODAY", "0.761", AMBER,
     "Frozen DINOv2-small · single seed · zero training · fully audited and reproducible from raw data "
     "to dashboard."),
    ("OURS — NEAR-TERM HEADROOM", "~0.83–0.90", BLUE,
     "Mask the letterbox padding, upgrade to DINOv2-base, run multi-seed. Engineering steps — "
     "not a redesign."),
    ("PUBLISHED SOTA", "0.95–0.96", MUT,
     "CSAD / SALAD: component segmentation, heavy training pipelines, far higher complexity budgets."),
]
for i, (head, num, col, body) in enumerate(cards19):
    cx = 0.6 + i * 4.20
    card(s, cx, 1.70, 3.95, 3.6)
    text(s, cx + 0.3, 1.95, 3.35, 0.35, [[{"t": head, "f": HEAD_F, "s": 12, "c": MUT, "b": True}]])
    text(s, cx + 0.3, 2.35, 3.35, 0.75, [[{"t": num, "f": HEAD_F, "s": 36, "c": col, "b": True}]])
    text(s, cx + 0.3, 3.25, 3.35, 1.85, [[{"t": body, "s": 12, "c": MUT}]], line=1.12)
card(s, 0.6, 5.62, 12.25, 1.4, fill=CARD2)
text(s, 0.92, 5.82, 11.7, 1.1, [
    [{"t": "We compete on rigor per unit of compute.", "f": HEAD_F, "s": 14.5, "c": AMBER, "b": True}],
    [{"t": "Each method is named for what it actually computes, evaluation is leakage-free, and the data "
           "is audited end to end. The remaining gap is understood and itemised — not mysterious.",
      "s": 12.5, "c": TXT}],
], line=1.12, space_after=6)
page(s, 19)
notes(s, "Pre-empt the obvious challenge ('SOTA is 0.95, you have 0.76'). The three-card framing converts "
         "it: we know exactly which engineering steps buy the next ~0.1, and we did not spend our budget "
         "on training — we spent it on auditability.")

# =========================================================================
# SLIDE 20 — Limitations & roadmap
# =========================================================================
s = slide()
kicker(s, "Limitations & roadmap")
title(s, "What we don't claim — and what comes next")
card(s, 0.6, 1.62, 6.0, 5.0)
text(s, 0.92, 1.84, 5.4, 0.4, [[{"t": "Known limitations", "f": HEAD_F, "s": 16, "c": RED, "b": True}]])
lims = [
    "Image-level metrics only — sPRO / pixel-level localization not yet reported.",
    "Letterbox padding still feeds the feature extractor (small, tracked bias).",
    "Single-seed results — no variance estimates yet.",
    "Global Image-Stat is a proxy, not the full EfficientAD it stands in for.",
    "Logical anomalies that require exact counting stay hard (pushpins logical 0.600).",
]
text(s, 0.92, 2.40, 5.4, 4.0,
     [[{"t": "—  ", "s": 12.3, "c": RED}, {"t": l, "s": 12.3, "c": MUT}] for l in lims],
     line=1.12, space_after=10)
card(s, 6.85, 1.62, 6.0, 5.0)
text(s, 7.17, 1.84, 5.4, 0.4, [[{"t": "Prioritised roadmap", "f": HEAD_F, "s": 16, "c": GREEN, "b": True}]])
roads = [
    ("P0", "Mask letterbox padding in feature space; re-run the full benchmark."),
    ("P0", "Backbone upgrade DINOv2-small → base (expected ≈ +0.05–0.08)."),
    ("P1", "Multi-seed runs; report mean ± std for every headline number."),
    ("P1", "Add sPRO and pixel-level localization to evaluation + dashboard."),
    ("P2", "Component-aware logical branch (CSAD-style) for the counting categories."),
]
for i, (tag, body) in enumerate(roads):
    py = 2.42 + i * 0.82
    chip(s, 7.17, py, 0.62, tag, color=AMBER if tag == "P0" else (BLUE if tag == "P1" else MUT), h=0.34, size=10.5)
    text(s, 7.95, py - 0.02, 4.6, 0.7, [[{"t": body, "s": 12.3, "c": MUT}]], line=1.05)
page(s, 20)
notes(s, "Symmetry sells it: every limitation on the left has a numbered fix on the right. P0 items are "
         "in-progress engineering, not aspirations.")

# =========================================================================
# SLIDE 21 — Conclusions
# =========================================================================
s = slide()
kicker(s, "Conclusions")
title(s, "Four things to take away")
concl = [
    ("1", "A leakage-free pipeline, end to end",
     "Read-only raw data, hashed splits with zero duplicates, validation-only calibration, held-out test. "
     "Every artifact is a CSV you can re-run."),
    ("2", "Fusion of complementary detectors works",
     "Four frozen-DINOv2 branches → 0.761 overall AUROC (0.731 logical / 0.804 structural), beating every "
     "single branch without training anything."),
    ("3", "Feature selection beats feature accumulation",
     "Leave-one-out analysis drops a duplicate and a harmful branch: the leaner 3-branch model reaches "
     "0.764 overall and 0.826 structural."),
    ("4", "Explainability is a deliverable, not a slide",
     "A self-contained offline dashboard shows score, threshold and evidence for every test decision — "
     "failures included."),
]
for i, (num, head, body) in enumerate(concl):
    cx = 0.6 + (i % 2) * 6.25
    cy = 1.75 + (i // 2) * 2.55
    card(s, cx, cy, 6.0, 2.3)
    text(s, cx + 0.3, cy + 0.25, 0.8, 0.6, [[{"t": num, "f": HEAD_F, "s": 26, "c": AMBER, "b": True}]])
    text(s, cx + 1.05, cy + 0.28, 4.65, 0.65, [[{"t": head, "f": HEAD_F, "s": 14.5, "c": TXT, "b": True}]], line=1.0)
    text(s, cx + 1.05, cy + 0.95, 4.65, 1.25, [[{"t": body, "s": 11.8, "c": MUT}]], line=1.08)
text(s, 0.6, 6.70, 12.2, 0.4, [[
    {"t": "Rigor, reproducibility and explainability — that is the contribution. The leaderboard number "
          "is the start, not the point.", "s": 12.5, "c": AMBER, "i": True}]])
page(s, 21)
notes(s, "Close with the four numbered cards in order — they recapitulate the agenda. End on the amber "
         "line; it is the thesis sentence of the whole project.")

# =========================================================================
# SLIDE 22 — Thank you
# =========================================================================
s = slide()
kicker(s, "Thank you", x=0.9, y=2.0)
text(s, 0.9, 2.35, 8.0, 1.1, [[{"t": "Questions welcome.", "f": HEAD_F, "s": 40, "c": TXT, "b": True}]])
text(s, 0.9, 3.6, 7.6, 0.9, [[
    {"t": "Everything in this talk regenerates from the repository — data audit → EDA → methods → "
          "results → dashboard → these slides.", "s": 14, "c": MUT}]], line=1.15)
card(s, 0.9, 4.75, 6.6, 0.85, fill=CARD2)
text(s, 1.2, 4.95, 6.0, 0.5, [[
    {"t": "Live demo:  ", "f": HEAD_F, "s": 14, "c": AMBER, "b": True},
    {"t": "09_dashboard/dashboard.html — fully offline", "f": "Consolas", "s": 13, "c": TXT}]])
for (img, px) in [("ex_breakfast_logical.png", 8.6), ("ex_pushpins_logical.png", 10.75)]:
    pic(s, f"{A}/{img}", px, 2.3, w=2.0, h=2.0)
for (img, px) in [("ex_juice_structural.png", 8.6), ("ex_splicing_logical.png", 10.75)]:
    pic(s, f"{A}/{img}", px, 4.45, w=2.0, h=2.0)
page(s, 22)
notes(s, "Offer the live dashboard for the Q&A screen. Likely questions: why not train the backbone "
         "(answer: reproducibility budget + frozen features already close most of the gap), and how the "
         "0.83–0.90 headroom estimate was derived (padding/backbone/multi-seed itemisation on slide 19).")

prs.save("Project_A_LOCO_Presentation_v2.pptx")
print(f"Saved Project_A_LOCO_Presentation_v2.pptx with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
