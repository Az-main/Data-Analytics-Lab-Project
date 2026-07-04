# QA: estimate rendered text size with real Windows fonts; flag overflow & collisions.
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont
import os

FONTS = {
    ("Trebuchet MS", False): r"C:\Windows\Fonts\trebuc.ttf",
    ("Trebuchet MS", True):  r"C:\Windows\Fonts\trebucbd.ttf",
    ("Calibri", False):      r"C:\Windows\Fonts\calibri.ttf",
    ("Calibri", True):       r"C:\Windows\Fonts\calibrib.ttf",
    ("Consolas", False):     r"C:\Windows\Fonts\consola.ttf",
    ("Consolas", True):      r"C:\Windows\Fonts\consolab.ttf",
}
_cache = {}
def font_for(name, bold, size_pt):
    key = (name, bold, round(size_pt * 2))
    if key not in _cache:
        path = FONTS.get((name, bold)) or FONTS.get((name, False)) or FONTS[("Calibri", bold)]
        _cache[key] = ImageFont.truetype(path, int(round(size_pt * 4)))  # 4x for precision
    return _cache[key]

def run_w(text, name, size, bold):
    f = font_for(name, bold, size)
    return f.getlength(text) / 4.0 / 72.0  # back to inches (pt measured at 4x)

def inch(v): return Emu(v).inches

prs = Presentation("Project_A_LOCO_Presentation_v2.pptx")
issues = []
for sidx, slide in enumerate(prs.slides, 1):
    shapes = []
    for sh in slide.shapes:
        if sh.left is None: continue
        shapes.append(sh)
    # compute estimated text bottoms
    boxes = []  # (shape, x, y, w, h, est_bottom, est_maxw)
    for sh in shapes:
        x, y, w, h = inch(sh.left), inch(sh.top), inch(sh.width), inch(sh.height)
        est_bottom = None
        if sh.has_text_frame and sh.text_frame.text.strip():
            total_h = 0.0
            maxw_needed = 0.0
            for p in sh.text_frame.paragraphs:
                runs = [(r.text, r.font.name or "Calibri", (r.font.size.pt if r.font.size else 13),
                         bool(r.font.bold)) for r in p.runs]
                if not runs:
                    continue
                size = max(r[2] for r in runs)
                ls = p.line_spacing if p.line_spacing else 1.0
                sa = (p.space_after.pt if p.space_after else 0) / 72.0
                # wrap simulation: split into words keeping run fonts approximately (use widest font of para)
                words = []
                for t, fn, sz, bd in runs:
                    for wd in t.split(" "):
                        words.append((wd, fn, sz, bd))
                space_w = run_w(" ", runs[0][1], size, runs[0][3])
                line_w, nlines = 0.0, 1
                avail = w if sh.text_frame.word_wrap is not False else 100
                for wd, fn, sz, bd in words:
                    ww = run_w(wd, fn, sz, bd)
                    add = ww if line_w == 0 else ww + space_w
                    if line_w + add > avail and line_w > 0:
                        nlines += 1
                        line_w = ww
                    else:
                        line_w += add
                    maxw_needed = max(maxw_needed, min(line_w, avail))
                line_h = size * 1.22 * ls / 72.0
                total_h += nlines * line_h + sa
            est_bottom = y + total_h
            boxes.append((sh, x, y, w, h, est_bottom, maxw_needed))
        else:
            boxes.append((sh, x, y, w, h, None, None))
    # checks
    cards = [(b[1], b[2], b[3], b[4]) for b in boxes
             if b[0].shape_type == 5 or (b[0].shape_type is not None and "ROUNDED" in str(b[0].shape_type))]
    # identify rounded rects via autoshape type
    cards = []
    for b in boxes:
        sh = b[0]
        try:
            if sh.shape_type == 1 and sh.adjustments and sh.height > Emu(0):
                from pptx.enum.shapes import MSO_SHAPE
                if sh.auto_shape_type in (MSO_SHAPE.ROUNDED_RECTANGLE,):
                    cards.append((b[1], b[2], b[3], b[4]))
        except Exception:
            pass
    for b in boxes:
        sh, x, y, w, h, est_b, maxw = b
        if est_b is None: continue
        snippet = sh.text_frame.text.replace("\n", " | ")[:60]
        # overflow past slide bottom
        if est_b > 7.45:
            issues.append(f"S{sidx}: text past slide bottom ({est_b:.2f}in): '{snippet}'")
        # overflow within its containing card
        for (cx, cy, cw, ch) in cards:
            if cx - 0.02 <= x <= cx + cw and cy - 0.02 <= y <= cy + ch and x + 0.1 < cx + cw:
                if est_b > cy + ch + 0.02:
                    issues.append(f"S{sidx}: text overflows card (card_bottom={cy+ch:.2f}, text_bottom={est_b:.2f}): '{snippet}'")
                break
        # severe overflow vs own declared box (>35% taller) — warning only
        if est_b > y + h * 1.6 and h > 0.25:
            issues.append(f"S{sidx}: WARN text much taller than its box (box_h={h:.2f}, est_h={est_b-y:.2f}): '{snippet}'")
        # single-word-too-wide check
        if maxw and maxw > w + 0.05:
            issues.append(f"S{sidx}: text wider than box (need {maxw:.2f} > {w:.2f}): '{snippet}'")
    # text-vs-picture collision (text starting above a picture but estimated to run into it)
    pics = [(b[1], b[2], b[3], b[4]) for b in boxes if b[0].shape_type == 13]
    for b in boxes:
        sh, x, y, w, h, est_b, _ = b
        if est_b is None: continue
        snippet = sh.text_frame.text.replace("\n", " | ")[:50]
        for (px, py, pw, ph) in pics:
            xo = max(x, px) < min(x + w, px + pw)   # x-overlap
            if xo and y < py and est_b > py + 0.05:
                issues.append(f"S{sidx}: text may run into picture below (text_bottom={est_b:.2f}, pic_top={py:.2f}): '{snippet}'")

print(f"{len(issues)} potential issues")
for i in issues: print(" -", i)
